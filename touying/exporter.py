import typst
from pptx import Presentation
from pptx.util import Cm
from pathlib import Path
from PIL import Image
from io import BytesIO
import json
import jinja2
import os
import re

FILE_PATH = os.path.dirname(os.path.realpath(__file__))


def to_html(
    input, root=None, font_paths=[], output=None, start_page=1, count=None, silent=False, sys_inputs={}
):
    if not silent:
        print(f"Compiling typst source file {input}...")
    
    images = typst.compile(input, root=root, font_paths=font_paths,
                           format="svg", sys_inputs=sys_inputs)
    if type(images) is not list:
        images = [images]

    # convert bytes to string
    images = [image.decode("utf-8") for image in images]
    # replace width="[0-9\.]+pt" height="[0-9\.]+pt" with width="100%" height="100%"
    images = [
        re.sub(
            r'width="([0-9\.]+)pt" height="([0-9\.]+)pt"',
            'width="100%" height="100%"',
            image,
        )
        for image in images
    ]

    # query <pdfpc-file> from typst file
    pdfpc = json.loads(
        typst.query(
            input, "<pdfpc-file>", root=root, font_paths=font_paths, field="value",
            sys_inputs=sys_inputs
        )
    )
    if len(pdfpc) > 0:
        pdfpc = pdfpc[0]
        idx2note = {
            page["idx"]: page["note"] for page in pdfpc["pages"] if "note" in page
        }
    else:
        pdfpc = None
        idx2note = {}

    if not silent:
        print("Creating presentation...")

    # create page iterator
    if count is None:
        count = len(images)
    page_iter = range(start_page - 1, start_page + count - 1)

    result = (
        jinja2.Environment(loader=jinja2.FileSystemLoader(FILE_PATH))
        .get_template("template.html.j2")
        .render(page_iter=page_iter, images=images, idx2note=idx2note, pdfpc=pdfpc)
    )

    # save to .html file
    if output is None:
        output = Path(input).with_suffix(".html")

    # create output directory if it doesn't exist
    output_dir = Path(output).parent
    if not output_dir.exists():
        output_dir.mkdir(parents=True)

    with open(output, "w", encoding="utf8") as f:
        f.write(result)


def to_pptx(
    input,
    output=None,
    root=None,
    font_paths=[],
    start_page=1,
    count=None,
    ppi=500,
    silent=False,
    sys_inputs={},
):

    if not silent:
        print(f"Compiling typst source file {input}...")

    images = typst.compile(
        input, root=root, font_paths=font_paths, format="png", ppi=ppi,
        sys_inputs=sys_inputs
    )
    if type(images) is not list:
        images = [images]

    # query <pdfpc-file> from typst file
    pdfpc = json.loads(
        typst.query(
            input, "<pdfpc-file>", root=root, font_paths=font_paths, field="value",
            sys_inputs=sys_inputs
        )
    )
    if len(pdfpc) > 0:
        pdfpc = pdfpc[0]
        idx2note = {
            page["idx"]: page["note"] for page in pdfpc["pages"] if "note" in page
        }
    else:
        pdfpc = None
        idx2note = {}

    if not silent:
        print("Creating presentation...")

    # create pptx presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # configure presentation aspect ratio
    page = Image.open(BytesIO(images[0]))
    aspect_ratio = page.width / page.height
    prs.slide_width = int(prs.slide_height * aspect_ratio)

    # create page iterator
    if count is None:
        count = len(images)
    page_iter = range(start_page - 1, start_page + count - 1)

    # iterate over slides
    for page_no in page_iter:
        image_file = BytesIO(images[page_no])

        # add a slide
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Cm(0)
        slide.shapes.add_picture(image_file, left, top, height=prs.slide_height)

        # add speaker notes
        if page_no in idx2note:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = idx2note[page_no]

    if output is None:
        output = Path(input).with_suffix(".pptx")

    # save presentation
    prs.save(output)

    # create output directory if it doesn't exist
    output_dir = Path(output).parent
    if not output_dir.exists():
        output_dir.mkdir(parents=True)

    if not silent:
        print(f"Presentation saved to {output}")


def to_pdf(input, output=None, root=None, font_paths=[], silent=False, sys_inputs={}):
    if not silent:
        print(f"Compiling typst source file {input}...")
    if output is None:
        output = Path(input).with_suffix(".pdf")
    # create output directory if it doesn't exist
    output_dir = Path(output).parent
    if not output_dir.exists():
        output_dir.mkdir(parents=True)
    typst.compile(input, output=output, root=root, font_paths=font_paths, format="pdf",
                  sys_inputs=sys_inputs)


def to_pdfpc(input, output=None, root=None, font_paths=[], silent=False, sys_inputs={}):
    if not silent:
        print(f"Compiling typst source file {input}...")

    # save to .pdfpc file
    if output is None:
        output = Path(input).with_suffix(".pdfpc")

    # create output directory if it doesn't exist
    output_dir = Path(output).parent
    if not output_dir.exists():
        output_dir.mkdir(parents=True)

    with open(output, "w", encoding="utf8") as f:
        f.write(
            typst.query(
                input,
                "<pdfpc-file>",
                root=root,
                font_paths=font_paths,
                field="value",
                one=True,
                sys_inputs=sys_inputs,
            )
        )


if __name__ == "__main__":
    # to_pptx("example.typ")
    to_html("example.typ")

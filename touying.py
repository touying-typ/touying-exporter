import typst
from pptx import Presentation
from pptx.util import Cm
from pathlib import Path
from PIL import Image
from io import BytesIO


def to_pptx(input, output=None, start_page=1, count=None, ppi=500, silent=False):

    if not silent:
        print(f"Compiling typst source file {input}...")

    images = typst.compile(input, format="png", ppi=ppi)

    if not silent:
        print("Creating presentation...")

    if type(images) is not list:
        images = [images]

    # create pptx presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    # configure presentation aspect ratio
    page = Image.open(BytesIO(images[0]))
    aspect_ratio = page.width / page.height
    prs.slide_width = int(prs.slide_height * aspect_ratio)

    # create page iterator
    if count is None:
        count = len(images)
    page_iter = range(start_page, start_page + count)

    # iterate over slides
    for page_no in page_iter:
        image_file = BytesIO(images[page_no - 1])

        # add a slide
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Cm(0)
        slide.shapes.add_picture(image_file, left, top, height=prs.slide_height)

    if output is None:
        output = Path(input).with_suffix(".pptx")

    # save presentation
    prs.save(output)

    if not silent:
        print(f"Presentation saved to {output}")


if __name__ == "__main__":
    to_pptx("example.typ")
